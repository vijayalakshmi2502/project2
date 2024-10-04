# app/main.py

import os
import re
import json
import logging
import time
import base64
import requests
from fastapi import FastAPI, Request, Form, UploadFile, File, status
from fastapi.responses import HTMLResponse, RedirectResponse
from fastapi.staticfiles import StaticFiles
from fastapi.templating import Jinja2Templates
import pandas as pd
from PyPDF2 import PdfFileReader
from docx import Document
from pptx import Presentation
from csv import reader
import xml.etree.ElementTree as ET
import azure.cognitiveservices.speech as speechsdk
from pydub import AudioSegment
from deep_translator import GoogleTranslator
from langdetect import detect
from azure.ai.textanalytics import TextAnalyticsClient
from azure.core.credentials import AzureKeyCredential
from azure.ai.translation.text import TextTranslationClient
from azure.cognitiveservices.vision.customvision.prediction import CustomVisionPredictionClient
from msrest.authentication import ApiKeyCredentials
import wave
import pyaudio
from fastapi.responses import JSONResponse


logging.basicConfig(
    level=logging.DEBUG,
    format='%(asctime)s - %(levelname)s - %(message)s',
    handlers=[
        logging.FileHandler("app.log"),
        logging.StreamHandler()
    ]
)


# Initialize FastAPI
app = FastAPI()

# Mount static files
app.mount("/static", StaticFiles(directory="static"), name="static")

# Initialize templates
templates = Jinja2Templates(directory="templates")

# Hard-coded Azure Credentials
DIRECT_LINE_SECRET="OjkTlOKRoRk.DX5CsW7aLKrJ50jDIXT3EvY0SuJ8OO7WADryBrIK5P4"
ENDPOINT = "https://p2ic-prediction.cognitiveservices.azure.com/"
PREDICTION_KEY = "64bfcfb8949d484fae93402fa55e1425"
PROJECT_ID = "c450da0e-a657-4c10-a1a4-f8ed35e8d40a"
PUBLISH_NAME = "DiseaseClassification"

# Create CustomVisionPredictionClient
credentials = ApiKeyCredentials(in_headers={"Prediction-key": PREDICTION_KEY})
predictor = CustomVisionPredictionClient(ENDPOINT, credentials)
# Azure credentials
AI_SERVICE_ENDPOINT = "https://p2languageservice.cognitiveservices.azure.com/"
AI_SERVICE_KEY = "741f912d0fb040d1afb1a3ca8b4f64d9"
TRANSLATION_ENDPOINT = "https://api.cognitive.microsofttranslator.com/"
TRANSLATION_KEY = "256e3786ca11469980ddae934f5944e5"
SPEECH_KEY = "872d1354c89e49158adb2b10b0c42493"
SPEECH_REGION = "eastus"

# Load all the datasets into a list of DataFrames
datasets = [
    pd.read_csv("D:/AI-Projects/FinalProject/P2/Lifestyle-Related Diseases.csv"),
    pd.read_csv("D:/AI-Projects/FinalProject/P2/disease_data.csv"),
    pd.read_csv("D:/AI-Projects/FinalProject/P2/Environmental Diseases.csv"),
    pd.read_csv("D:/AI-Projects/FinalProject/P2/Idiopathic.csv"),
    pd.read_csv("D:/AI-Projects/FinalProject/P2/Neoplastic Diseases.csv"),
    pd.read_csv("D:/AI-Projects/FinalProject/P2/non-infectious diseases_data.csv"),
    pd.read_csv("D:/AI-Projects/FinalProject/P2/Nutritional Diseases.csv"),
    pd.read_csv("D:/AI-Projects/FinalProject/P2/Psychiatric and Neurological Disorders.csv"),
    pd.read_csv("D:\AI-Projects\FinalProject\P2\Rare Diseases.csv")
]

# Combine all datasets into one DataFrame for easier searching
combined_df = pd.concat(datasets, ignore_index=True)

# Language code to full name dictionary
lang_code_to_full = {
    'en': 'English',
    'fr': 'French',
    'es': 'Spanish',
    'de': 'German',
    'it': 'Italian',
    'pt': 'Portuguese',
    'zh': 'Chinese',
    'ja': 'Japanese',
    'ko': 'Korean',
    'ar': 'Arabic',
    'ru': 'Russian',
    'hi': 'Hindi',
    'bn': 'Bengali',
    'pa': 'Punjabi',
    'jv': 'Javanese',
    'ml': 'Malayalam',
    'tr': 'Turkish',
    'vi': 'Vietnamese',
    'th': 'Thai',
    'sv': 'Swedish',
    'da': 'Danish',
    'no': 'Norwegian',
    'fi': 'Finnish',
    'cs': 'Czech',
    'pl': 'Polish',
    'uk': 'Ukrainian',
    'el': 'Greek',
    'he': 'Hebrew',
    'te': 'Telugu',
    'ur': 'Urdu',
    'kn': 'Kannada',
    'ta': 'Tamil'
}

# Initialize Azure Clients
def authenticate_azure_translation():
    credential = AzureKeyCredential(TRANSLATION_KEY)
    client = TextTranslationClient(endpoint=TRANSLATION_ENDPOINT, credential=credential)
    return client

def authenticate_azure_nlp():
    credential = AzureKeyCredential(AI_SERVICE_KEY)
    client = TextAnalyticsClient(endpoint=AI_SERVICE_ENDPOINT, credential=credential)
    return client

def authenticate_custom_vision():
    credentials = ApiKeyCredentials(in_headers={"Prediction-key": PREDICTION_KEY})
    predictor = CustomVisionPredictionClient(ENDPOINT, credentials)
    return predictor

predictor = authenticate_custom_vision()

# Utility Functions
def clean_text(text):
    text = re.sub(r'\W+', ' ', text.lower())
    text = ' '.join(dict.fromkeys(text.split()))  # Remove duplicate words
    return text

def detect_language(text):
    client = authenticate_azure_nlp()
    try:
        response = client.detect_language(documents=[text])
        primary_language = response[0].primary_language.iso6391_name
        return primary_language
    except Exception as e:
        logging.error(f"Error detecting language: {e}")
        return None

def translate_to_english(text):
    try:
        user_lang = detect(text)
        logging.debug(f"Detected Language: {lang_code_to_full.get(user_lang, user_lang)}")
        if user_lang != 'en':
            text = GoogleTranslator(source=user_lang, target='en').translate(text)
    except Exception as e:
        logging.error(f"Error translating to English: {e}")
    return text

def extract_key_phrases(text):
    client = authenticate_azure_nlp()
    try:
        response = client.extract_key_phrases(documents=[text])[0]
        return response.key_phrases
    except Exception as e:
        logging.error(f"Error extracting key phrases: {e}")
        return []

def find_disease_by_symptoms(key_phrases, combined_df):
    key_phrases_cleaned = [clean_text(phrase) for phrase in key_phrases]
    matched_diseases = combined_df[combined_df['Symptoms'].apply(
        lambda x: any(phrase in clean_text(x) for phrase in key_phrases_cleaned)
    )]
    return matched_diseases

def prioritize_combined(matched_diseases, key_phrases):
    severity_priority = {'Mild': 1, 'Moderate': 2, 'Severe': 3}
    matched_diseases = matched_diseases.copy()
    matched_diseases.loc[:, 'Severity Priority'] = matched_diseases['Severity Level'].map(severity_priority)
    
    def count_matches(symptoms):
        cleaned_symptoms = clean_text(symptoms)
        return sum(phrase in cleaned_symptoms for phrase in key_phrases)
    
    matched_diseases.loc[:, 'Match Count'] = matched_diseases['Symptoms'].apply(count_matches)
    prioritized_diseases = matched_diseases.sort_values(by=['Match Count', 'Severity Priority'], ascending=[False, False])
    return prioritized_diseases
    
def get_disease_info(disease, language):
    disease_info = {
        "Disease Name": disease['Disease Name'],
        "Severity Level": disease['Severity Level'],
        "Symptoms": disease['Symptoms'],
        "Recommended Medications": disease['Recommended Medications'],
        "Required Food": disease['Required Food'],
        "Safety Precautions": disease['Safety Precautions'],
        "Recommended Doctor": disease['Recommended Doctor'],
        "Treatment Plan": disease['Treatment Plan'],
        "Follow-Up Recommendations": disease['Follow-Up Recommendations'],
        "Patient Education": disease['Patient Education'],
        "Recovery Time": disease['Recovery Time']
    }
    
    if language != 'en':
        translator = GoogleTranslator(source='en', target=language)
        translated_info = {key: translator.translate(value) for key, value in disease_info.items()}
        return translated_info
    
    return disease_info



def detect_language_and_sentiment(text):
    client = authenticate_azure_nlp()
    
    # Detect Language
    try:
        detected_language = client.detect_language(documents=[text])[0]
        detected_lang_code = detected_language.primary_language.iso6391_name
        logging.debug(f"Detected Language: {lang_code_to_full.get(detected_lang_code, detected_lang_code)}")
    except Exception as e:
        logging.error(f"Error detecting language: {e}")
        detected_lang_code = 'en'
    
    # Sentiment Analysis
    try:
        sentiment_response = client.analyze_sentiment(documents=[text])[0]
        sentiment = sentiment_response.sentiment
        logging.debug(f"Sentiment: {sentiment}")
    except Exception as e:
        logging.error(f"Error analyzing sentiment: {e}")
        sentiment = 'neutral'
    
    return sentiment, detected_lang_code

def extract_text_from_file(file_path):
    # Extract text based on file type
    try:
        if file_path.endswith('.txt'):
            with open(file_path, 'r', encoding='utf-8') as file:
                return file.read()
        elif file_path.endswith('.pdf'):
            with open(file_path, 'rb') as file:
                reader = PdfFileReader(file)
                text = ""
                for page_num in range(reader.numPages):
                    text += reader.getPage(page_num).extract_text()
                return text
        elif file_path.endswith('.docx'):
            doc = Document(file_path)
            text = ""
            for para in doc.paragraphs:
                text += para.text
            return text
        elif file_path.endswith('.pptx'):
            presentation = Presentation(file_path)
            text = ""
            for slide in presentation.slides:
                for shape in slide.shapes:
                    if hasattr(shape, 'text'):
                        text += shape.text + "\n"
            return text
        elif file_path.endswith('.xlsx'):
            df = pd.read_excel(file_path)
            text = ""
            for column in df.columns:
                text += df[column].astype(str).str.cat(sep='\n')
            return text
        elif file_path.endswith('.csv'):
            text = ""
            with open(file_path, 'r', encoding='utf-8') as file:
                csv_reader = reader(file)
                for row in csv_reader:
                    text += ' '.join(row) + "\n"
            return text
        elif file_path.endswith('.json'):
            with open(file_path, 'r', encoding='utf-8') as file:
                data = json.load(file)
                return json.dumps(data, indent=2)
        elif file_path.endswith('.xml'):
            tree = ET.parse(file_path)
            root = tree.getroot()
            return ET.tostring(root, encoding='unicode')
    except Exception as e:
        logging.error(f"Error extracting text from file: {e}")
        return ""
    return ""

def record_live_audio(file_path, duration=5):
    """Record live audio and save it as a WAV file."""
    try:
        audio = pyaudio.PyAudio()
        stream = audio.open(format=pyaudio.paInt16, channels=1, rate=44100, input=True, frames_per_buffer=1024)
        
        logging.debug("Recording...")
        frames = []
        
        for _ in range(int(44100 / 1024 * duration)):
            data = stream.read(1024)
            frames.append(data)
        
        logging.debug("Recording finished.")
        stream.stop_stream()
        stream.close()
        audio.terminate()
        
        with wave.open(file_path, 'wb') as wf:
            wf.setnchannels(1)
            wf.setsampwidth(audio.get_sample_size(pyaudio.paInt16))
            wf.setframerate(44100)
            wf.writeframes(b''.join(frames))
        logging.debug(f"Recorded audio saved to {file_path}")
        return file_path
    except Exception as e:
        logging.error(f"Error recording audio: {e}")
        return None


def transcribe_audio(file_path):
    speech_config = speechsdk.SpeechConfig(subscription=SPEECH_KEY, region=SPEECH_REGION)
    audio_config = speechsdk.AudioConfig(filename=file_path)
    speech_recognizer = speechsdk.SpeechRecognizer(speech_config=speech_config, audio_config=audio_config)
    
    all_text = []
    recognition_finished = False

    def recognized_callback(evt):
        if evt.result.reason == speechsdk.ResultReason.RecognizedSpeech:
            logging.debug(f"Recognized: {evt.result.text}")
            all_text.append(evt.result.text)
        elif evt.result.reason == speechsdk.ResultReason.NoMatch:
            logging.debug("No speech could be recognized.")

    def session_stopped_callback(evt):
        nonlocal recognition_finished
        recognition_finished = True
        logging.debug("Recognition session stopped.")

    speech_recognizer.recognized.connect(recognized_callback)
    speech_recognizer.session_stopped.connect(session_stopped_callback)
    
    logging.debug("Transcribing...")
    speech_recognizer.start_continuous_recognition()

    try:
        while not recognition_finished:
            time.sleep(1)
    except KeyboardInterrupt:
        speech_recognizer.stop_continuous_recognition()
    
    speech_recognizer.stop_continuous_recognition()

    return " ".join(all_text)

def classify_image(image_path):
    with open(image_path, "rb") as image_data:
        results = predictor.classify_image(PROJECT_ID, PUBLISH_NAME, image_data.read())
    
    logging.info(f"Predictions for image: {image_path}")
    for prediction in results.predictions:
        logging.info(f"Tag: {prediction.tag_name}, Probability: {prediction.probability:.2f}")

    best_prediction = max(results.predictions, key=lambda p: p.probability)
    return best_prediction.tag_name if best_prediction.probability > 0 else None

# Routes
@app.get("/about", response_class=HTMLResponse)
async def about(request: Request):
    return templates.TemplateResponse("about.html", {"request": request})

@app.get("/contact", response_class=HTMLResponse)
async def contact(request: Request):
    return templates.TemplateResponse("contact.html", {"request": request})


@app.get("/", response_class=HTMLResponse)
async def read_index(request: Request):
    return templates.TemplateResponse("index.html", {"request": request, "languages": lang_code_to_full})


@app.post("/process-symptoms", response_class=HTMLResponse)
async def process_symptoms(request: Request, symptoms: str = Form(...), preferred_language: str = Form(...)):
    try:
        print(symptoms)
        translated_text = translate_to_english(symptoms)
        print(translated_text)
        sentiment, detected_language = detect_language_and_sentiment(translated_text)
        
        if sentiment == 'positive':
            result = "The input statement is positive. No disease identification needed."
            print(result)
            return templates.TemplateResponse("response.html", {"request": request, "message": result})
        
        key_phrases = extract_key_phrases(translated_text)
        print(f"\nExtracted Key Phrases (Symptoms): {key_phrases}")
        matched_diseases = find_disease_by_symptoms(key_phrases, combined_df)
        
        if not matched_diseases.empty:
            prioritized_diseases = prioritize_combined(matched_diseases, key_phrases)
            top_disease = prioritized_diseases.iloc[0]
            
            # Ensure preferred_language is valid
            if preferred_language not in lang_code_to_full:
                preferred_language = 'en'  # Default to English if invalid
            
            disease_info =  get_disease_info(top_disease, preferred_language)
            return templates.TemplateResponse("response.html", {"request": request, "predicted_disease": disease_info})
        else:
            result = "No matching diseases found based on the provided symptoms."
            return templates.TemplateResponse("response.html", {"request": request, "message": result})
    
    except Exception as e:
        # Log the error or return an error message
        error_message = f"An error occurred: {str(e)}"
        return templates.TemplateResponse("response.html", {"request": request, "predicted_disease": error_message})

@app.post("/upload-file", response_class=HTMLResponse)
async def upload_file(request: Request, file: UploadFile = File(...), preferred_language: str = Form(...)):
    upload_dir = "uploads/files/"
    os.makedirs(upload_dir, exist_ok=True)
    file_path = os.path.join(upload_dir, file.filename)
    
    with open(file_path, "wb") as buffer:
        buffer.write(await file.read())
    
    file_text = extract_text_from_file(file_path)
    print(file_text)
    translated_text = translate_to_english(file_text)
    sentiment, detected_language = detect_language_and_sentiment(translated_text)
    
    if sentiment == 'positive':
        result = "The input statement is positive. No disease identification needed."
        print(result)
        return templates.TemplateResponse("response.html", {"request": request, "message": result})
    
    key_phrases = extract_key_phrases(translated_text)
    print(f"\nExtracted Key Phrases (Symptoms): {key_phrases}")
    matched_diseases = find_disease_by_symptoms(key_phrases, combined_df)
    
    if not matched_diseases.empty:
        prioritized_diseases = prioritize_combined(matched_diseases, key_phrases)
        top_disease = prioritized_diseases.iloc[0]
        
        if preferred_language not in lang_code_to_full:
            preferred_language = 'en'
        
        disease_info =  get_disease_info(top_disease, preferred_language)
        return templates.TemplateResponse("response.html", {"request": request, "predicted_disease": disease_info})
    else:
        result = "No matching diseases found based on the provided symptoms."
        return templates.TemplateResponse("response.html", {"request": request, "message": result})

@app.post("/upload-audio", response_class=HTMLResponse)
async def upload_audio(request: Request, audio: UploadFile = File(...), preferred_language: str = Form(...)):
    upload_dir = "uploads/audio/"
    os.makedirs(upload_dir, exist_ok=True)
    file_path = os.path.join(upload_dir, audio.filename)
    
    with open(file_path, "wb") as buffer:
        buffer.write(await audio.read())
    
    transcription = transcribe_audio(file_path)
    print(transcription)
    translated_text = translate_to_english(transcription)
    sentiment, detected_language = detect_language_and_sentiment(translated_text)
    
    if sentiment == 'positive':
        result = "The input statement is positive. No disease identification needed."
        print(result)
        return templates.TemplateResponse("response.html", {"request": request, "message": result})
    
    key_phrases = extract_key_phrases(translated_text)
    print(f"\nExtracted Key Phrases (Symptoms): {key_phrases}")
    matched_diseases = find_disease_by_symptoms(key_phrases, combined_df)
    
    if not matched_diseases.empty:
        prioritized_diseases = prioritize_combined(matched_diseases, key_phrases)
        top_disease = prioritized_diseases.iloc[0]
        
        if preferred_language not in lang_code_to_full:
            preferred_language = 'en'
        
        disease_info =  get_disease_info(top_disease, preferred_language)
        return templates.TemplateResponse("response.html", {"request": request, "predicted_disease": disease_info})
    else:
        result = "No matching diseases found based on the provided symptoms."
        return templates.TemplateResponse("response.html", {"request": request, "message": result})

@app.post("/record-audio", response_class=HTMLResponse)
async def record_audio(
    request: Request,
    duration: int = Form(...),
    audio_data: str = Form(...),
    preferred_language: str = Form(...),
):
    upload_dir = "uploads/audio/"
    os.makedirs(upload_dir, exist_ok=True)
    webm_file_path = os.path.join(upload_dir, "live_recorded_audio.webm")
    wav_file_path = os.path.join(upload_dir, "live_recorded_audio.wav")
    print(f"Saving audio to: {webm_file_path}")

    # Decode base64 audio data and save as WebM file
    try:
        if not audio_data:
            raise ValueError("No audio data received.")

        audio_bytes = base64.b64decode(audio_data)
        with open(webm_file_path, "wb") as audio_file:
            audio_file.write(audio_bytes)
        print("Audio file saved successfully as WebM.")
    except Exception as e:
        logging.error(f"Error saving audio data: {e}")
        result = "Failed to process the recorded audio."
        return templates.TemplateResponse("response.html", {"request": request, "message": result})

    # Convert WebM to WAV using pydub
    try:
        audio = AudioSegment.from_file(webm_file_path, format="webm")
        audio.export(wav_file_path, format="wav")
        print("Audio file converted to WAV successfully.")
    except Exception as e:
        logging.error(f"Error converting audio file to WAV: {e}")
        result = "Failed to convert audio to WAV format."
        return templates.TemplateResponse("response.html", {"request": request, "message": result})

    # Proceed with transcription and further processing
    try:
        transcription = transcribe_audio(wav_file_path)
        print(f"Transcription result: {transcription}")
    except Exception as e:
        logging.error(f"Error during transcription: {e}")
        transcription = "Transcription failed."

    # Proceed only if transcription succeeded
    if transcription and transcription != "Transcription failed.":
        translated_text = translate_to_english(transcription)
        sentiment, detected_language = detect_language_and_sentiment(translated_text)

        if sentiment == 'positive':
            result = "The input statement is positive. No disease identification needed."
            print(result)
            return templates.TemplateResponse("response.html", {"request": request, "message": result})

        key_phrases = extract_key_phrases(translated_text)
        print(f"Extracted Key Phrases (Symptoms): {key_phrases}")
        matched_diseases = find_disease_by_symptoms(key_phrases, combined_df)

        if not matched_diseases.empty:
            prioritized_diseases = prioritize_combined(matched_diseases, key_phrases)
            top_disease = prioritized_diseases.iloc[0]

            if preferred_language not in lang_code_to_full:
                preferred_language = 'en'

            disease_info =  get_disease_info(top_disease, preferred_language)
            return templates.TemplateResponse("response.html", {"request": request, "predicted_disease": disease_info})
        else:
            result = "No matching diseases found based on the provided symptoms."
            return templates.TemplateResponse("response.html", {"request": request, "message": result})
    else:
        result = "Transcription failed or returned no meaningful data."
        return templates.TemplateResponse("response.html", {"request": request, "message": result})


@app.post("/upload-image", response_class=HTMLResponse)
async def upload_image(request: Request, image: UploadFile = File(...), preferred_language: str = Form(...)):
    upload_dir = "uploads/images/"
    os.makedirs(upload_dir, exist_ok=True)
    file_path = os.path.join(upload_dir, image.filename)
    
    with open(file_path, "wb") as buffer:
        buffer.write(await image.read())
    
    disease_tag = classify_image(file_path)
    if disease_tag:
        matched_diseases = combined_df[combined_df['Disease Name'].str.lower() == disease_tag.lower()]
        
        if not matched_diseases.empty:
            prioritized_diseases = prioritize_combined(matched_diseases, [])
            top_disease = prioritized_diseases.iloc[0]
            
            if preferred_language not in lang_code_to_full:
                preferred_language = 'en'
            
            disease_info = get_disease_info(top_disease, preferred_language)
            return templates.TemplateResponse("response.html", {"request": request, "predicted_disease": disease_info})
        else:
            result = "No matching diseases found for the classified tag."
            return templates.TemplateResponse("response.html", {"request": request, "message": result})
    else:
        result = "No disease could be classified from the image."
        return templates.TemplateResponse("response.html", {"request": request, "message": result})

# === Azure Health Bot Integration ===

@app.get("/healthbot", response_class=HTMLResponse)
async def healthbot(request: Request):
    """
    Serve the Azure Health Bot page.
    """
    try:
        return templates.TemplateResponse("hbot.html", {"request": request})
    except Exception as e:
        logging.error(f"Error serving Health Bot page: {e}")
        return HTMLResponse(content="An error occurred while loading the Health Bot.", status_code=500)

@app.get("/directline/token", response_class=JSONResponse)
async def get_directline_token():
    """
    Generate a Direct Line token for Azure Health Bot.
    """
    token_url = "https://directline.botframework.com/v3/directline/tokens/generate"
    headers = {
        "Authorization": f"Bearer {DIRECT_LINE_SECRET}",
        "Content-Type": "application/json"
    }

    try:
        response = requests.post(token_url, headers=headers)
        response.raise_for_status()
        return response.json()  # Contains the token
    except requests.exceptions.HTTPError as http_err:
        logging.error(f"HTTP error occurred while generating token: {http_err}")
        return JSONResponse(content={"error": "Failed to generate token"}, status_code=response.status_code)
    except Exception as e:
        logging.error(f"Error occurred while generating token: {e}")
        return JSONResponse(content={"error": "An unexpected error occurred"}, status_code=500)



if __name__ == '__main__':
    import uvicorn
    uvicorn.run(app, host="0.0.0.0", port=5002, log_level="debug")

#uvicorn p2app:app --host 0.0.0.0 --port 5002
#http://localhost:5002/